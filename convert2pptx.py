import sys
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import os

def pdf_to_ppt(pdf_path, pptx_path):
    print("Converting PDF to images...")
    # This converts PDF pages into a list of PIL Image objects
    images = convert_from_path(pdf_path)
    
    prs = Presentation()
    # Set slide dimensions to a standard 16:9 ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print("Adding images to slides...")
    for i, image in enumerate(images):
        # Save temporary image
        image_path = f"temp_page_{i}.png"
        image.save(image_path, 'PNG')
        
        # Add blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add image to slide
        left = top = Inches(0)
        slide.shapes.add_picture(image_path, left, top, width=prs.slide_width, height=prs.slide_height)
        
        # Remove temporary image
        os.remove(image_path)
        
    prs.save(pptx_path)
    print(f"PPTX saved to {pptx_path}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python convert2ppt.py <input.pdf> <output.pptx>")
    else:
        pdf_to_ppt(sys.argv[1], sys.argv[2])